import xlrd
from pptx import Presentation
from pptx.util import Inches, Pt
import time
import glob
import os
from pptx.dml.color import RGBColor
import numpy as np



start = time.time()
prs = Presentation('C:\\**sample.pptx')
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
shapes = slide.shapes
title = slide.shapes.title
title.text = "OPCO name"
img_path = 'C:**\\Romania.png' #any image need to insert
pic = slide.shapes.add_picture(img_path, Inches(6.5), Inches(0.5),
                               width=Inches(3), height=Inches(2))

#genearet color blocks similar to color plate but with choosen colors
def colorlist(value,start_value,start_color,end_value,end_color):
    output = []
    value = int(np.round(value,0))
    r1, g1, b1 = RGBColor.from_string(start_color)
    r2, g2, b2 = RGBColor.from_string(end_color)
    rdelta, gdelta, bdelta = (r2-r1)/(end_value-start_value), (g2-g1)/(end_value-start_value), (b2-b1)/(end_value-start_value)
    for step in range(0,end_value - start_value):
        r1 += rdelta
        g1 += gdelta
        b1 += bdelta
        output.append((int(r1), int(g1), int(b1)))
    color = output[value-start_value]
    return color
#personilizing color applying rules
def ColorScaleRule(value,start_value,start_color, mid_value,mid_color,end_color,end_value):
    value = int(np.round(value,0))

    if value in range(start_value,mid_value):
        return colorlist(value,start_value,start_color,mid_value,mid_color)   
    elif value in range(mid_value,end_value):
        return colorlist(value,mid_value,mid_color,end_value,end_color)
    else:
        return RGBColor.from_string(end_color)
#choosing some colors range to srat with        
def CPURAMColorScaleRule(value):
    start_value = 0 
    start_color = '98fb98'
    mid_value = 50
    mid_color = '7cfc00'
    end_color = '228b22'
    end_value = 80 
    return ColorScaleRule(value,start_value,start_color, mid_value,mid_color,end_color,end_value)      
#choosing some colors range to srat with  - for a diff slide
def OtherColorScaleRule(sheetname,value):
    start_color = '00ff7f'
    mid_color = 'ffff00'
    end_color = 'ffff00'
    if sheetname in ('Contention & Ready Time','CPU Co-Stop'):
        start_value = 0
        mid_value = 5
        end_value = 5
        color = ColorScaleRule(value,start_value,start_color, mid_value,mid_color,end_color,end_value)
    elif sheetname == 'Latency' :
        start_value = 0
        mid_value = 15
        end_value = 15
        color = ColorScaleRule(value,start_value,start_color, mid_value,mid_color,end_color,end_value)
    
    elif sheetname in ('Packet Drop','VMs at Risk'):
        start_value = 0
        mid_value = 1
        end_value = 1
        if value == 0:
            color = RGBColor.from_string('00ff7f')
        else:
            color = RGBColor.from_string('ffff00')
    return color 
#getting RGB codes
def getRGBcolor(tup):
    r= int(tup[0])
    b= int(tup[1])
    g= int(tup[2])
    return RGBColor(r,b,g)
def main():
    for filename in glob.glob(os.path.join("C:,"GR_NVI-MET-DC02_AVG__SUMMARY.xlsx")):
        print(filename)
        file_location = filename
        try:            
            workbook = xlrd.open_workbook(file_location)
            sheetnames = workbook.sheet_names()
            for sheetname in sheetnames:
                sheet = workbook.sheet_by_name(sheetname)
                print("sheet:", sheetname)
                rows = sheet.nrows
                cols = sheet.ncols
                c = cols
                r = rows
                if c > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    shapes = slide.shapes
                    title = slide.shapes.title
                    title.text = "OPCO name"
                    left = Inches(0.1)
                    right = Inches(0.1)
                    top = Inches(1.0)
                    width = Inches(4.0)
                    height = Inches(0.5)
                    num = 10.0/c
                    
                    table = shapes.add_table(rows, cols, left, top, width, right).table
                    fill = slide.shape.fill
                    fill.solid()
                    fill.fore_color.brightness = -0.25
                    
                    for i in range(0,c):
                        table.columns[i].width = Inches(num)
                    for i in range(0,r):
                        for e in range(0,c):
                            table.cell(i,e).text = str(sheet.cell_value(i,e))
                            cell = table.rows[i].cells[e]
                            paragraph = cell.text_frame.paragraphs[0]
                            paragraph.font.size = Pt(11)
                            if sheetname in ('CPU','Memory'):
                                if i!=0 and e in (2,3,5):                                    
                                    fill = cell.fill #fill the legend as well
                                    fill.solid()
                                    tup = CPURAMColorScaleRule(sheet.cell_value(i,e))
                                    fill.fore_color.rgb = getRGBcolor(tup)
                            elif sheetname in ('Latency & Packet Drop'):
                                if i!=0 and e in (1,2,3):
                                    fill = cell.fill #fill the legend as well
                                    fill.solid()
                                    sheetname = 'Latency'
                                    tup = OtherColorScaleRule(sheetname ,(sheet.cell_value(i,e)))
                                    fill.fore_color.rgb = getRGBcolor(tup)
                                elif i!=0 and e in (4,5,6):
                                    fill = cell.fill #fill the legend as well
                                    fill.solid()
                                    sheetname = 'Packet Drop'
                                    tup = OtherColorScaleRule(sheetname ,(sheet.cell_value(i,e)))
                                    fill.fore_color.rgb = getRGBcolor(tup)
                            elif sheetname in('Contention & Ready Time','CPU Co-Stop'):
                                if i!=0 and e!=0 :
                                    fill = cell.fill #fill the legend as well
                                    fill.solid()
                                    tup = OtherColorScaleRule(sheetname ,(sheet.cell_value(i,e)))
                                    fill.fore_color.rgb = getRGBcolor(tup)
                            elif sheetname == 'VMs at Risk':
                                if i!=0 and not e in (0,1) :
                                    fill = cell.fill #fill the legend as well
                                    fill.solid()
                                    tup = OtherColorScaleRule(sheetname ,(sheet.cell_value(i,e)))
                                    fill.fore_color.rgb = getRGBcolor(tup)
                                    
        except Exception as e:
            print("Error!" )
            print(e)
            
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    title = slide.shapes.title
    title.text = "Heatmaps"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
    left = Inches(0.7)
    right = Inches(0.1)
    top = Inches(1.7)
    width = Inches(4.0)
    rows=3
    cols = 7
    table = shapes.add_table(rows, cols, left, top, width, right).table
    column_header = ['CPU Demand','Memory Demand','CPU Contention','CPU Ready Time', 'Storage Latency', 'NW Packet Drop']
    for i in range(0,cols):
        table.columns[i].width = Inches(2.5)
        cell = table.cell(0, i)
        cell.text
        cell.text = column_header[i]  
        
    prs.save(os.path.join("C:**//OPCO_OUTPUT",'powerpointfile1.pptx'))
    end = time.time()
    print(end - start)

    
if __name__ == '__main__':
    main()